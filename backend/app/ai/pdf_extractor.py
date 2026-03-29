"""
PadaiSathi AI — Document Text Extractor
Place at:  backend/app/ai/pdf_extractor.py

Supports PDF, PPTX, and TXT files.
Tries PyMuPDF first (better quality), falls back to PyPDF2 for PDFs.
"""

import os
import re


def extract_text_from_pdf(filepath: str) -> str:
    """
    Extract and clean text from a PDF, PPTX, or TXT file.
    Returns cleaned string ready for summarization.
    """
    if not os.path.exists(filepath):
        raise FileNotFoundError(f"File not found: {filepath}")

    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".txt":
        raw = _extract_txt(filepath)
    elif ext == ".pptx":
        raw = _extract_pptx(filepath)
    else:
        raw = _try_pymupdf(filepath) or _try_pypdf2(filepath)

    if not raw or not raw.strip():
        raise ValueError(
            "Could not extract text from this file. "
            "It may be empty or in an unsupported format."
        )

    return _clean(raw)


def _extract_txt(filepath: str) -> str:
    """Extract text from a plain .txt file."""
    with open(filepath, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    print(f"[Extractor] TXT: {len(text)} chars")
    return text


def _extract_pptx(filepath: str) -> str:
    """Extract text from a .pptx file (all slides, all shapes)."""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = " ".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
        text = "\n".join(parts)
        print(f"[Extractor] PPTX: {len(text)} chars")
        return text
    except ImportError:
        raise ImportError("python-pptx is required for PPTX support. Run: pip install python-pptx")
    except Exception as e:
        print(f"[Extractor] PPTX error: {e}")
        return ""


def _try_pymupdf(filepath: str) -> str:
    """Best quality extractor. pip install pymupdf"""
    try:
        import fitz
        parts = []
        with fitz.open(filepath) as doc:
            for page in doc:
                parts.append(page.get_text("text"))
        text = "\n".join(parts)
        if text.strip():
            print(f"[PDFExtractor] PyMuPDF: {len(text)} chars")
            return text
    except ImportError:
        pass
    except Exception as e:
        print(f"[PDFExtractor] PyMuPDF error: {e}")
    return ""


def _try_pypdf2(filepath: str) -> str:
    """Fallback — uses your existing PyPDF2 install."""
    try:
        import PyPDF2
        text = ""
        with open(filepath, "rb") as f:
            reader = PyPDF2.PdfReader(f)
            for page in reader.pages:
                text += (page.extract_text() or "") + "\n"
        if text.strip():
            print(f"[PDFExtractor] PyPDF2: {len(text)} chars")
            return text
    except ImportError:
        pass
    except Exception as e:
        print(f"[PDFExtractor] PyPDF2 error: {e}")
    return ""


def _clean(text: str) -> str:
    """Remove noise from extracted PDF text."""
    # Normalise unicode punctuation
    for old, new in [('\u2013', '-'), ('\u2014', '-'), ('\u2018', "'"),
                     ('\u2019', "'"), ('\u201c', '"'), ('\u201d', '"')]:
        text = text.replace(old, new)

    # Remove lines that are just page numbers
    lines = [ln for ln in text.splitlines() if not re.match(r'^\s*\d+\s*$', ln)]
    text = '\n'.join(lines)

    # Collapse excessive blank lines and spaces
    text = re.sub(r'\n{3,}', '\n\n', text)
    text = re.sub(r'[ \t]{2,}', ' ', text)

    return text.strip()